"""提取样例PPT和PDF的文本内容，用于长期规划参考。"""
import sys
from pathlib import Path

def extract_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"\n===== Slide {i} =====")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        lines.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip("| "):
                        lines.append(row_text)
    return "\n".join(lines)

def extract_pdf(path: str) -> str:
    import pymupdf
    doc = pymupdf.open(path)
    lines = []
    for i, page in enumerate(doc, 1):
        lines.append(f"\n===== Page {i} =====")
        text = page.get_text()
        if text.strip():
            lines.append(text)
    return "\n".join(lines)

if __name__ == "__main__":
    pptx_path = r"D:\AAA-Study\work\样例\LLM-KG-Carto一种基于大语言模型和知识图谱的自动制图框架-廖成-2021202050034.pptx"
    pdf_path = r"D:\AAA-Study\work\样例\艾廷华--DeepSeek驱动下的地图生成.pdf"

    out_dir = Path(r"D:\AAA-Study\work\github\carto-agent\.preview\samples")
    out_dir.mkdir(parents=True, exist_ok=True)

    print("Extracting PPTX...")
    pptx_text = extract_pptx(pptx_path)
    (out_dir / "llm_kg_carto_pptx.txt").write_text(pptx_text, encoding="utf-8")
    print(f"PPTX: {len(pptx_text)} chars -> {out_dir / 'llm_kg_carto_pptx.txt'}")

    print("Extracting PDF...")
    pdf_text = extract_pdf(pdf_path)
    (out_dir / "deepseek_map_pdf.txt").write_text(pdf_text, encoding="utf-8")
    print(f"PDF: {len(pdf_text)} chars -> {out_dir / 'deepseek_map_pdf.txt'}")
